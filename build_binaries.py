#!/usr/bin/env python3
"""
Build script to create binaries for DocuGenius
Supports both macOS and Windows binary creation
"""

import os
import sys
import subprocess
import shutil
from pathlib import Path
import tempfile

def run_command(cmd, capture_output=True):
    """Run a command and return success status"""
    try:
        if capture_output:
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
            return result.returncode == 0, result.stdout, result.stderr
        else:
            result = subprocess.run(cmd, shell=True)
            return result.returncode == 0, "", ""
    except Exception as e:
        return False, "", str(e)

def create_cli_source():
    """Create the CLI source code with integrated image extraction"""
    cli_source = '''#!/usr/bin/env python3
"""
DocuGenius CLI - Document to Markdown Converter with Image Extraction
A standalone document converter for DocuGenius VS Code extension
"""

import sys
import os
import argparse
from pathlib import Path
import json
import re
import hashlib
from typing import List, Dict, Tuple, Optional

def convert_text_file(file_path):
    """Convert text-based files (just read and return content)"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return content
    except UnicodeDecodeError:
        # Try with different encodings
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                return content
            except UnicodeDecodeError:
                continue
        return f"# {Path(file_path).name}\\n\\nError: Could not decode file content."

def convert_json_file(file_path):
    """Convert JSON file to formatted markdown"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        markdown = f"# {Path(file_path).name}\\n\\n"
        markdown += "```json\\n"
        markdown += json.dumps(data, indent=2, ensure_ascii=False)
        markdown += "\\n```\\n"
        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting JSON: {str(e)}"

def convert_csv_file(file_path):
    """Convert CSV file to markdown table"""
    try:
        import csv
        markdown = f"# {Path(file_path).name}\\n\\n"

        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return markdown + "Empty CSV file."

        # Header row
        if rows:
            markdown += "| " + " | ".join(rows[0]) + " |\\n"
            markdown += "| " + " | ".join(["---"] * len(rows[0])) + " |\\n"

            # Data rows
            for row in rows[1:]:
                # Pad row to match header length
                padded_row = row + [""] * (len(rows[0]) - len(row))
                markdown += "| " + " | ".join(padded_row[:len(rows[0])]) + " |\\n"

        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting CSV: {str(e)}"

def convert_xml_file(file_path):
    """Convert XML file to formatted markdown"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        markdown = f"# {Path(file_path).name}\\n\\n"
        markdown += "```xml\\n"
        markdown += content
        markdown += "\\n```\\n"
        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting XML: {str(e)}"

def convert_docx_file(file_path):
    """Convert DOCX file using python-docx"""
    try:
        from docx import Document

        doc = Document(file_path)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Handle different paragraph styles
                style_name = paragraph.style.name.lower() if paragraph.style else ""

                if "heading 1" in style_name:
                    markdown += f"# {text}\\n\\n"
                elif "heading 2" in style_name:
                    markdown += f"## {text}\\n\\n"
                elif "heading 3" in style_name:
                    markdown += f"### {text}\\n\\n"
                elif "heading 4" in style_name:
                    markdown += f"#### {text}\\n\\n"
                elif "heading 5" in style_name:
                    markdown += f"##### {text}\\n\\n"
                elif "heading 6" in style_name:
                    markdown += f"###### {text}\\n\\n"
                else:
                    markdown += f"{text}\\n\\n"

        # Extract tables
        for table in doc.tables:
            markdown += "\\n"
            for i, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace('\\n', ' ')
                    row_data.append(cell_text)

                if i == 0:
                    # Header row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
                    markdown += "| " + " | ".join(["---"] * len(row_data)) + " |\\n"
                else:
                    # Data row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
            markdown += "\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: python-docx library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting DOCX: {str(e)}"

def convert_excel_file(file_path):
    """Convert Excel file using openpyxl"""
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, data_only=True)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]

            markdown += f"## {sheet_name}\\n\\n"

            # Get all rows with data
            rows = list(worksheet.iter_rows(values_only=True))
            if not rows:
                markdown += "*Empty sheet*\\n\\n"
                continue

            # Filter out completely empty rows
            non_empty_rows = []
            for row in rows:
                if any(cell is not None and str(cell).strip() for cell in row):
                    non_empty_rows.append(row)

            if not non_empty_rows:
                markdown += "*No data found*\\n\\n"
                continue

            # Find the maximum number of columns with data
            max_cols = max(len([cell for cell in row if cell is not None]) for row in non_empty_rows)

            # Create markdown table
            for i, row in enumerate(non_empty_rows):
                # Convert row to strings and pad to max_cols
                row_data = []
                for j in range(max_cols):
                    if j < len(row) and row[j] is not None:
                        row_data.append(str(row[j]).strip())
                    else:
                        row_data.append("")

                if i == 0:
                    # Header row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
                    markdown += "| " + " | ".join(["---"] * len(row_data)) + " |\\n"
                else:
                    # Data row
                    markdown += "| " + " | ".join(row_data) + " |\\n"

            markdown += "\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: openpyxl library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting Excel: {str(e)}"

def convert_pptx_file(file_path):
    """Convert PowerPoint file using python-pptx"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        for i, slide in enumerate(prs.slides, 1):
            markdown += f"## Slide {i}\\n\\n"

            # Extract text from all shapes in the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                for text in slide_text:
                    # Split by lines and format appropriately
                    lines = text.split('\\n')
                    for line in lines:
                        line = line.strip()
                        if line:
                            markdown += f"{line}\\n\\n"
            else:
                markdown += "*No text content found*\\n\\n"

            markdown += "---\\n\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: python-pptx library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting PowerPoint: {str(e)}"

def convert_pdf_file(file_path):
    """Convert PDF file using pdfplumber"""
    try:
        import pdfplumber

        file_name = Path(file_path).name
        markdown = f"# {file_name}\\n\\n"

        with pdfplumber.open(file_path) as pdf:
            markdown += f"**Total Pages:** {len(pdf.pages)}\\n\\n"

            for i, page in enumerate(pdf.pages, 1):
                markdown += f"## Page {i}\\n\\n"

                try:
                    text = page.extract_text()
                    if text and text.strip():
                        # Clean up the extracted text
                        lines = text.split('\\n')
                        cleaned_lines = []
                        for line in lines:
                            line = line.strip()
                            if line:
                                cleaned_lines.append(line)

                        if cleaned_lines:
                            markdown += '\\n\\n'.join(cleaned_lines) + "\\n\\n"
                        else:
                            markdown += "*No text content found on this page*\\n\\n"
                    else:
                        markdown += "*No text content found on this page*\\n\\n"

                except Exception as page_error:
                    markdown += f"*Error extracting text from page {i}: {str(page_error)}*\\n\\n"

                markdown += "---\\n\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: pdfplumber library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting PDF: {str(e)}"

def extract_images_from_pdf(file_path, output_dir, min_image_size=50):
    """PDF image extraction not supported in lightweight mode"""
    return [], "PDF image extraction is not supported in lightweight mode (pdfplumber does not support image extraction)"

def convert_document_file(file_path, extract_images=True):
    """Convert document files using native Python libraries with optional image extraction"""
    file_name = Path(file_path).name
    file_ext = Path(file_path).suffix.lower()

    try:
        # First, convert the document content
        if file_ext in ['.docx']:
            content = convert_docx_file(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            content = convert_excel_file(file_path)
        elif file_ext in ['.pptx']:
            content = convert_pptx_file(file_path)
        elif file_ext == '.pdf':
            content = convert_pdf_file(file_path)
        else:
            # Fallback for unsupported formats
            content = f"# {file_name}\\n\\n"
            content += f"**Document Type:** {file_ext.upper()} file\\n\\n"
            content += "This file type is not yet supported for full conversion.\\n\\n"
            content += f"- **File:** {file_name}\\n"
            content += f"- **Size:** {os.path.getsize(file_path)} bytes\\n\\n"
            return content

        # If image extraction is enabled and we have a PDF
        if extract_images and file_ext == '.pdf':
            # Note about image extraction limitation
            content += "\\n\\n<!-- Note: PDF image extraction is not supported in lightweight mode (using pdfplumber) -->\\n"

        return content

    except Exception as e:
        # Error handling - return basic info with error message
        content = f"# {file_name}\\n\\n"
        content += f"**Error converting {file_ext.upper()} file**\\n\\n"
        content += f"Error: {str(e)}\\n\\n"
        content += f"- **File:** {file_name}\\n"
        content += f"- **Size:** {os.path.getsize(file_path)} bytes\\n"
        return content


def main():
    if len(sys.argv) < 2:
        print("DocuGenius CLI - Document to Markdown Converter", file=sys.stderr)
        print("Usage: docugenius-cli <file> [extract_images]", file=sys.stderr)
        print("", file=sys.stderr)
        print("Arguments:", file=sys.stderr)
        print("  file           : Path to document file", file=sys.stderr)
        print("  extract_images : true/false to enable/disable image extraction for DOCX/PPTX/XLSX (default: true)", file=sys.stderr)
        print("", file=sys.stderr)
        print("Supported formats:", file=sys.stderr)
        print("  - Text files: .txt, .md, .markdown", file=sys.stderr)
        print("  - Data files: .json, .csv, .xml, .html", file=sys.stderr)
        print("  - Documents: .docx, .xlsx, .pptx (with image extraction), .pdf (text only)", file=sys.stderr)
        print("", file=sys.stderr)
        print("Features:", file=sys.stderr)
        print("  - Converts documents to Markdown format", file=sys.stderr)
        print("  - High-quality text extraction from PDF files (using pdfplumber)", file=sys.stderr)
        print("  - Lightweight and cross-platform consistent", file=sys.stderr)
        print("  - Fast installation and execution", file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    extract_images = True

    if len(sys.argv) > 2:
        extract_images = sys.argv[2].lower() not in ['false', 'no', '0']

    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    file_ext = Path(file_path).suffix.lower()

    try:
        if file_ext in ['.txt', '.md', '.markdown']:
            content = convert_text_file(file_path)
        elif file_ext == '.json':
            content = convert_json_file(file_path)
        elif file_ext == '.csv':
            content = convert_csv_file(file_path)
        elif file_ext in ['.xml', '.html', '.htm']:
            content = convert_xml_file(file_path)
        elif file_ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf']:
            content = convert_document_file(file_path, extract_images)
        else:
            # Default to text file handling for unknown extensions
            content = convert_text_file(file_path)

        print(content)

    except Exception as e:
        print(f"Error processing file: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
'''
    return cli_source

def create_darwin_binary():
    """Create macOS Universal Binary using PyInstaller (supports both Intel and Apple Silicon)"""
    print("🔨 Building DocuGenius macOS Universal Binary")
    print("=" * 50)

    # Detect current architecture
    import platform
    current_arch = platform.machine()
    print(f"🏗️  Current system architecture: {current_arch}")
    if current_arch == "arm64":
        print("   (Apple Silicon - ARM64)")
    elif current_arch == "x86_64":
        print("   (Intel - x86_64)")
    else:
        print(f"   (Unknown architecture: {current_arch})")
    
    print("🎯 Target: Universal Binary (Intel + Apple Silicon)")

    # Create temporary CLI source file
    with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
        f.write(create_cli_source())
        cli_file = f.name

    try:
        # Create virtual environment for building
        env_dir = "build_env_darwin"
        if os.path.exists(env_dir):
            shutil.rmtree(env_dir)

        print(f"📦 Creating build environment: {env_dir}")
        success, _, _ = run_command(f"python3 -m venv {env_dir}")
        if not success:
            print("❌ Failed to create virtual environment")
            return False

        # Install PyInstaller and document processing libraries
        print("📥 Installing PyInstaller and document libraries...")
        install_cmd = f"source {env_dir}/bin/activate && pip install pyinstaller python-docx python-pptx openpyxl pdfplumber"
        success, _, _ = run_command(install_cmd)

        if not success:
            print("❌ Failed to install required libraries")
            return False

        # Build the executable with universal binary support
        print("🔨 Building Universal Binary executable...")
        # Use --target-arch=universal2 for universal binary support
        build_cmd = f"source {env_dir}/bin/activate && python -m PyInstaller --onefile --name docugenius-cli --target-arch=universal2 --strip --optimize=2 {cli_file}"

        success, stdout, stderr = run_command(build_cmd, capture_output=False)

        if not success:
            print("⚠️  Universal binary build failed, falling back to current architecture...")
            # Fallback to current architecture if universal build fails
            build_cmd = f"source {env_dir}/bin/activate && python -m PyInstaller --onefile --name docugenius-cli --strip --optimize=2 {cli_file}"
            success, stdout, stderr = run_command(build_cmd, capture_output=False)
            
            if not success:
                print("❌ Failed to build executable")
                return False

        # Check if the executable was created
        exe_path = "dist/docugenius-cli"

        if not os.path.exists(exe_path):
            print("❌ Executable not found after build")
            return False

        # Create the bin/darwin directory if it doesn't exist
        darwin_dir = Path("bin/darwin")
        darwin_dir.mkdir(parents=True, exist_ok=True)

        # Copy the executable to the bin directory
        target_path = darwin_dir / "docugenius-cli"

        shutil.copy2(exe_path, target_path)
        os.chmod(target_path, 0o755)

        # Check if it's a universal binary
        try:
            result = subprocess.run(["lipo", "-info", str(target_path)], capture_output=True, text=True)
            if "Architectures in the fat file" in result.stdout:
                print(f"✅ Universal Binary created: {target_path}")
                print(f"🏗️  Supported architectures: {result.stdout.split(':')[-1].strip()}")
            else:
                arch_info = result.stdout.strip().split()[-1] if result.stdout else "unknown"
                print(f"✅ Single architecture binary created: {target_path}")
                print(f"🏗️  Architecture: {arch_info}")
                if current_arch == "x86_64" and arch_info == "x86_64":
                    print("ℹ️  Note: This binary will work on Apple Silicon via Rosetta 2")
                elif current_arch == "arm64" and arch_info == "arm64":
                    print("ℹ️  Note: This binary is optimized for Apple Silicon")
        except Exception as e:
            print(f"⚠️  Could not verify binary architecture: {e}")

        print(f"📊 File size: {os.path.getsize(target_path) / (1024*1024):.1f} MB")

        # Clean up build artifacts
        cleanup_dirs = ['build', 'dist', env_dir]
        for dir_name in cleanup_dirs:
            if os.path.exists(dir_name):
                shutil.rmtree(dir_name)

        print("🧹 Cleaned up build artifacts")
        return True

    except Exception as e:
        print(f"❌ Failed to create binary: {e}")
        return False
    finally:
        # Clean up temporary file
        if os.path.exists(cli_file):
            os.unlink(cli_file)

def create_windows_batch():
    """Create Windows batch file"""
    print("🔨 Creating DocuGenius Windows Batch File")
    print("=" * 40)

    # Create the bin/win32 directory if it doesn't exist
    win32_dir = Path("bin/win32")
    win32_dir.mkdir(parents=True, exist_ok=True)

    # Get the CLI source and modify it for Windows batch
    cli_source = create_cli_source()

    # Convert CLI source to Windows batch format
    # Replace Python source with Windows batch commands that create and run a temp Python script
    python_lines = cli_source.split('\n')[3:]  # Skip shebang and docstring start

    batch_content = '''@echo off
REM DocuGenius CLI for Windows
setlocal enabledelayedexpansion

if "%~1"=="" (
    echo DocuGenius CLI - Document to Markdown Converter
    echo Usage: docugenius-cli ^<file^>
    echo.
    echo Supported formats:
    echo   - Text files: .txt, .md, .markdown
    echo   - Data files: .json, .csv, .xml, .html
    echo   - Documents: .docx, .xlsx, .pptx, .pdf
    exit /b 1
)

if not exist "%~1" (
    echo Error: File not found: %~1
    exit /b 1
)

python --version >nul 2>&1
if errorlevel 1 (
    echo Error: Python is not installed or not in PATH
    echo Please install Python from https://python.org
    echo Note: You may need to install: pip install python-docx python-pptx openpyxl pdfplumber
    exit /b 1
)

REM Create temporary Python script
set TEMP_SCRIPT=%TEMP%\\docugenius_temp_%RANDOM%.py

(
'''

    # Add Python code to batch file, escaping special characters
    for line in python_lines:
        if line.strip():
            # Escape special batch characters
            escaped_line = line.replace('^', '^^').replace('&', '^&').replace('<', '^<').replace('>', '^>').replace('|', '^|')
            escaped_line = escaped_line.replace('(', '^(').replace(')', '^)')
            batch_content += f"echo {escaped_line}\n"
        else:
            batch_content += "echo.\n"

    batch_content += ''') > "%TEMP_SCRIPT%"

REM Run the Python script
python "%TEMP_SCRIPT%" "%~1"
set RESULT=%ERRORLEVEL%

REM Clean up temporary script
del "%TEMP_SCRIPT%" >nul 2>&1

exit /b %RESULT%
'''

    target_path = win32_dir / "docugenius-cli.bat"

    try:
        with open(target_path, 'w') as f:
            f.write(batch_content)

        print(f"✅ Windows batch file created: {target_path}")
        print(f"📊 File size: {os.path.getsize(target_path)} bytes")
        return True

    except Exception as e:
        print(f"❌ Failed to create batch file: {e}")
        return False

def main():
    print("🚀 DocuGenius Binary Builder")
    print("=" * 50)

    platform = sys.platform

    if len(sys.argv) > 1:
        target = sys.argv[1].lower()
    else:
        target = "all"

    success = True

    if target in ["all", "darwin", "macos"]:
        if platform == "darwin" or target != "all":
            success &= create_darwin_binary()
        else:
            print("⚠️  Skipping macOS binary (not on macOS)")

    if target in ["all", "windows", "win32"]:
        success &= create_windows_batch()

    if success:
        print("\n🎉 Binary build completed successfully!")
        print("\n📁 Generated files:")
        if os.path.exists("bin/darwin/docugenius-cli"):
            size = os.path.getsize("bin/darwin/docugenius-cli") / (1024*1024)
            print(f"   - bin/darwin/docugenius-cli ({size:.1f} MB)")
        if os.path.exists("bin/win32/docugenius-cli.bat"):
            size = os.path.getsize("bin/win32/docugenius-cli.bat")
            print(f"   - bin/win32/docugenius-cli.bat ({size} bytes)")
    else:
        print("\n💥 Build failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()
