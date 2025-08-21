# -*- coding: utf-8 -*-
"""
DocuGenius Image Extraction Module
Extracts images from PDF, DOCX, and PPTX files and organizes them in structured folders.
"""

import sys
import os
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Tuple, Optional

class ImageExtractor:
    """Main class for extracting images from documents"""

    def __init__(self, document_path: str, output_dir: str = None, markdown_dir: str = None, min_image_size: int = 50):
        self.document_path = Path(document_path)
        self.document_name = self.document_path.stem
        self.document_ext = self.document_path.suffix.lower()

        # Set output directory - use DocuGenius/images/{document_name}/ structure
        if output_dir:
            self.output_dir = Path(output_dir) / self.document_name
        else:
            # Default to DocuGenius/images folder in same directory as document
            self.output_dir = self.document_path.parent / "DocuGenius" / "images" / self.document_name

        self.extracted_images = []
        self.image_counter = 1
        self.min_image_size = min_image_size

        # Store the original document directory for relative path calculation
        self.original_dir = self.document_path.parent

        # Set markdown directory for relative path calculation
        if markdown_dir:
            self.markdown_dir = Path(markdown_dir)
        else:
            # Default to DocuGenius directory for markdown files
            self.markdown_dir = self.original_dir / "DocuGenius"

    def extract_images(self) -> Dict:
        """Extract images from document and return metadata"""
        try:
            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            if self.document_ext == '.pdf':
                return self._extract_from_pdf()
            elif self.document_ext == '.docx':
                return self._extract_from_docx()
            elif self.document_ext == '.pptx':
                return self._extract_from_pptx()
            elif self.document_ext == '.xlsx':
                return self._extract_from_xlsx()
            else:
                return {
                    'success': False,
                    'error': f'Unsupported document type: {self.document_ext}',
                    'images': []
                }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images: {str(e)}',
                'images': []
            }
    
    def _extract_from_pdf(self) -> Dict:
        """PDF image extraction is not supported in lightweight mode"""
        return {
            'success': False,
            'error': 'PDF image extraction is not supported in lightweight mode. DocuGenius uses pdfplumber for PDF text extraction only to keep dependencies small (0.8MB vs 45MB).',
            'images': [],
            'note': 'PDF files will be converted to text-only Markdown. Images in PDFs are not extracted to maintain lightweight dependencies.'
        }



    def _extract_from_docx(self) -> Dict:
        """Extract images from DOCX files"""
        try:
            import docx
            from docx.document import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            return {
                'success': False,
                'error': 'DOCX image extraction requires python-docx (pip install python-docx)',
                'images': []
            }

        try:
            doc = docx.Document(str(self.document_path))
            images_extracted = []

            # Extract images from document relationships
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    # Get image data
                    image_part = rel.target_part
                    image_data = image_part.blob

                    # Determine file extension from content type
                    content_type = image_part.content_type
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        img_ext = 'jpg'
                    elif 'png' in content_type:
                        img_ext = 'png'
                    elif 'gif' in content_type:
                        img_ext = 'gif'
                    elif 'bmp' in content_type:
                        img_ext = 'bmp'
                    else:
                        img_ext = 'png'  # Default fallback

                    # Generate unique filename
                    img_filename = self._generate_image_filename(
                        f"docx_img_{len(images_extracted) + 1}",
                        img_ext
                    )
                    img_path = self.output_dir / img_filename

                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_data)

                    # Add to extracted images list
                    image_info = {
                        'filename': img_filename,
                        'path': str(img_path),
                        'relative_path': f"images/{self.document_name}/{img_filename}",
                        'format': img_ext.upper(),
                        'size_bytes': len(image_data),
                        'source': 'docx_relationship'
                    }
                    images_extracted.append(image_info)
                    self.extracted_images.append(image_info)

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from DOCX: {str(e)}',
                'images': []
            }

    def _extract_from_pptx(self) -> Dict:
        """Extract images from PPTX files"""
        try:
            import pptx
            from pptx import Presentation
        except ImportError:
            return {
                'success': False,
                'error': 'PPTX image extraction requires python-pptx (pip install python-pptx)',
                'images': []
            }

        try:
            prs = Presentation(str(self.document_path))
            images_extracted = []

            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        # Get image data
                        image = shape.image
                        image_data = image.blob

                        # Determine file extension from content type
                        content_type = image.content_type
                        if 'jpeg' in content_type or 'jpg' in content_type:
                            img_ext = 'jpg'
                        elif 'png' in content_type:
                            img_ext = 'png'
                        elif 'gif' in content_type:
                            img_ext = 'gif'
                        elif 'bmp' in content_type:
                            img_ext = 'bmp'
                        else:
                            img_ext = 'png'  # Default fallback

                        # Generate unique filename
                        img_filename = self._generate_image_filename(
                            f"slide_{slide_num + 1}_img_{len([s for s in slide.shapes if hasattr(s, 'image')]) + 1}",
                            img_ext
                        )
                        img_path = self.output_dir / img_filename

                        # Save image
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_data)

                        # Add to extracted images list
                        image_info = {
                            'filename': img_filename,
                            'path': str(img_path),
                            'relative_path': f"images/{self.document_name}/{img_filename}",
                            'slide': slide_num + 1,
                            'format': img_ext.upper(),
                            'size_bytes': len(image_data),
                            'source': 'pptx_shape'
                        }
                        images_extracted.append(image_info)
                        self.extracted_images.append(image_info)

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from PPTX: {str(e)}',
                'images': []
            }

    def _generate_image_filename(self, base_name: str, extension: str) -> str:
        """Generate a unique filename for an image, avoiding collisions"""
        # Clean the base name
        base_name = "".join(c for c in base_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        base_name = base_name.replace(' ', '_')

        filename = f"{base_name}.{extension}"
        counter = 1

        # Check for collisions and add counter if needed
        while (self.output_dir / filename).exists():
            filename = f"{base_name}_{counter}.{extension}"
            counter += 1

        return filename

    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate MD5 hash of a file for duplicate detection"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def generate_markdown_references(self, images: List[Dict], mode: str = "simple") -> str:
        """Generate markdown image references for extracted images"""
        if not images:
            return ""

        markdown_lines = []

        if mode == "grouped":
            # Group images by page/slide
            grouped_images = {}
            for img in images:
                key = img.get('page', img.get('slide', 'unknown'))
                if key not in grouped_images:
                    grouped_images[key] = []
                grouped_images[key].append(img)

            markdown_lines.append("## Extracted Images")
            markdown_lines.append("")

            for key in sorted(grouped_images.keys()):
                if isinstance(key, int):
                    if 'page' in grouped_images[key][0]:
                        markdown_lines.append(f"### Page {key}")
                    else:
                        markdown_lines.append(f"### Slide {key}")
                else:
                    markdown_lines.append(f"### {key}")
                markdown_lines.append("")

                for img in grouped_images[key]:
                    alt_text = self._generate_alt_text(img)
                    markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
                    markdown_lines.append("")

        elif mode == "inline":
            # Just return the image references without headers
            for img in images:
                alt_text = self._generate_alt_text(img)
                markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
                markdown_lines.append("")

        else:  # simple mode
            markdown_lines.append("## Extracted Images")
            markdown_lines.append("")

            for img in images:
                alt_text = self._generate_alt_text(img)
                markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
                markdown_lines.append("")

        return "\n".join(markdown_lines)

    def _generate_alt_text(self, img: Dict) -> str:
        """Generate appropriate alt text for an image"""
        alt_text = "Image"

        if 'page' in img:
            alt_text += f" from page {img['page']}"
        elif 'slide' in img:
            alt_text += f" from slide {img['slide']}"
        elif 'sheet' in img:
            alt_text += f" from sheet {img['sheet']}"

        # Add dimensions if available
        if 'width' in img and 'height' in img:
            alt_text += f" ({img['width']}x{img['height']})"

        return alt_text

    def _calculate_relative_path(self, image_path: Path) -> str:
        """Calculate the correct relative path from markdown file to image"""
        try:
            # Use the instance markdown_dir for relative path calculation
            markdown_dir = self.markdown_dir

            # Calculate relative path from markdown file to image
            relative_path = os.path.relpath(str(image_path), str(markdown_dir))

            # Ensure forward slashes for markdown compatibility
            relative_path = relative_path.replace('\\', '/')

            return relative_path
        except Exception as e:
            # Fallback to simple relative path
            return f"images/{self.document_name}/{image_path.name}"

    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate MD5 hash of a file for duplicate detection"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def get_image_references_by_page(self, images: List[Dict]) -> Dict[int, List[str]]:
        """Get image references organized by page number for easy insertion"""
        page_images = {}

        for img in images:
            page = img.get('page', img.get('slide', 0))
            if page not in page_images:
                page_images[page] = []

            alt_text = f"Image"
            if 'page' in img:
                alt_text += f" from page {img['page']}"
            elif 'slide' in img:
                alt_text += f" from slide {img['slide']}"

            page_images[page].append(f"![{alt_text}]({img['relative_path']})")

        return page_images

    def get_simple_image_list(self, images: List[Dict]) -> List[str]:
        """Get a simple list of image markdown references"""
        image_refs = []

        for img in images:
            alt_text = f"Image"
            if 'page' in img:
                alt_text += f" from page {img['page']}"
            elif 'slide' in img:
                alt_text += f" from slide {img['slide']}"

            image_refs.append(f"![{alt_text}]({img['relative_path']})")

        return image_refs

    def extract_document_content_with_images(self) -> Dict:
        """Extract document content and intelligently insert images at their original positions"""
        try:
            if self.document_ext == '.pdf':
                return self._extract_pdf_content_with_images()
            elif self.document_ext == '.docx':
                return self._extract_docx_content_with_images()
            elif self.document_ext == '.pptx':
                return self._extract_pptx_content_with_images()
            elif self.document_ext == '.xlsx':
                return self._extract_xlsx_content_with_images()
            else:
                return {
                    'success': False,
                    'error': f'Unsupported document type for content extraction: {self.document_ext}',
                    'markdown_content': '',
                    'images': []
                }
        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_pdf_content_with_images(self) -> Dict:
        """PDF content with images is not supported in lightweight mode"""
        return {
            'success': False,
            'error': 'PDF image extraction is not supported in lightweight mode. DocuGenius uses pdfplumber for PDF text extraction only to keep dependencies small.',
            'markdown_content': '',
            'images': [],
            'note': 'Use simple PDF text conversion instead. Images in PDFs are not extracted to maintain lightweight dependencies.'
        }







    def _extract_docx_content_with_images(self) -> Dict:
        """Extract DOCX content and insert images at their original positions"""
        try:
            import docx
            from docx.document import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            return {
                'success': False,
                'error': 'DOCX content extraction requires python-docx (pip install python-docx)',
                'markdown_content': '',
                'images': []
            }

        try:
            doc = docx.Document(str(self.document_path))
            markdown_lines = []
            all_images = []

            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            # Extract images first
            image_counter = 1
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # Get image data
                        image_part = rel.target_part
                        image_data = image_part.blob

                        # Determine file extension from content type
                        content_type = image_part.content_type
                        if 'jpeg' in content_type or 'jpg' in content_type:
                            img_ext = 'jpg'
                        elif 'png' in content_type:
                            img_ext = 'png'
                        elif 'gif' in content_type:
                            img_ext = 'gif'
                        elif 'bmp' in content_type:
                            img_ext = 'bmp'
                        else:
                            img_ext = 'png'  # Default fallback

                        # Generate unique filename
                        img_filename = self._generate_image_filename(
                            f"docx_img_{image_counter}",
                            img_ext
                        )
                        img_path = self.output_dir / img_filename

                        # Save image
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_data)

                        # Create image info
                        image_info = {
                            'filename': img_filename,
                            'path': str(img_path),
                            'relative_path': self._calculate_relative_path(img_path),
                            'format': img_ext.upper(),
                            'size_bytes': len(image_data),
                            'source': 'docx_image'
                        }
                        all_images.append(image_info)
                        image_counter += 1

                    except Exception as img_error:
                        print(f"Warning: Could not extract image {image_counter}: {img_error}")
                        continue

            # Process document content
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    # Paragraph
                    paragraph = Paragraph(element, doc)
                    text = paragraph.text.strip()
                    if text:
                        markdown_lines.append(text)
                        markdown_lines.append("")

                    # Check if paragraph contains images (simplified approach)
                    # In a full implementation, we would track image positions more precisely
                    if len(all_images) > 0 and "image" in text.lower():
                        # Insert first available image here (simplified logic)
                        for img_info in all_images:
                            if img_info not in [img for img in all_images if img.get('inserted')]:
                                alt_text = f"Image from document ({img_info.get('format', 'Unknown')})"
                                markdown_lines.append(f"![{alt_text}]({img_info['relative_path']})")
                                markdown_lines.append("")
                                img_info['inserted'] = True
                                break

                elif isinstance(element, CT_Tbl):
                    # Table
                    table = Table(element, doc)
                    for i, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        if i == 0:  # Header row
                            markdown_lines.append("| " + " | ".join(row_data) + " |")
                            markdown_lines.append("| " + " | ".join(["---"] * len(row_data)) + " |")
                        else:
                            markdown_lines.append("| " + " | ".join(row_data) + " |")
                    markdown_lines.append("")

            # Insert any remaining images at the end
            for img_info in all_images:
                if not img_info.get('inserted'):
                    alt_text = f"Image from document ({img_info.get('format', 'Unknown')})"
                    markdown_lines.append(f"![{alt_text}]({img_info['relative_path']})")
                    markdown_lines.append("")

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(all_images),
                'images': all_images,
                'markdown_content': '\n'.join(markdown_lines)
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting DOCX content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_pptx_content_with_images(self) -> Dict:
        """Extract PPTX content and insert images at their original positions"""
        # This is a simplified version - full implementation would be more complex
        return {
            'success': False,
            'error': 'PPTX content extraction with images not fully implemented in this version',
            'markdown_content': '',
            'images': []
        }

    def _extract_xlsx_content_with_images(self) -> Dict:
        """Extract XLSX content and insert images at their original positions"""
        # This is a simplified version - full implementation would be more complex
        return {
            'success': False,
            'error': 'XLSX content extraction with images not fully implemented in this version',
            'markdown_content': '',
            'images': []
        }

    def _extract_from_xlsx(self) -> Dict:
        """Extract images from XLSX files"""
        try:
            import openpyxl
            from openpyxl.drawing.image import Image as OpenpyxlImage
        except ImportError:
            return {
                'success': False,
                'error': 'XLSX image extraction requires openpyxl (pip install openpyxl)',
                'images': []
            }

        try:
            workbook = openpyxl.load_workbook(str(self.document_path))
            images_extracted = []

            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]

                # Check if worksheet has images
                if hasattr(worksheet, '_images') and worksheet._images:
                    for img_index, image in enumerate(worksheet._images):
                        try:
                            # Get image data
                            image_data = image.ref.getvalue()

                            # Determine file extension from image format
                            img_ext = 'png'  # Default
                            if hasattr(image, 'format') and image.format:
                                img_ext = image.format.lower()
                            elif hasattr(image.ref, 'format') and image.ref.format:
                                img_ext = image.ref.format.lower()

                            # Generate unique filename
                            img_filename = self._generate_image_filename(
                                f"sheet_{sheet_name}_img_{img_index + 1}",
                                img_ext
                            )
                            img_path = self.output_dir / img_filename

                            # Save image
                            with open(img_path, "wb") as img_file:
                                img_file.write(image_data)

                            # Add to extracted images list
                            image_info = {
                                'filename': img_filename,
                                'path': str(img_path),
                                'relative_path': self._calculate_relative_path(img_path),
                                'sheet': sheet_name,
                                'format': img_ext.upper(),
                                'size_bytes': len(image_data),
                                'source': 'xlsx_image'
                            }
                            images_extracted.append(image_info)
                            self.extracted_images.append(image_info)

                        except Exception as img_error:
                            # Skip this image if there's an error
                            print(f"Warning: Could not extract image {img_index + 1} from sheet {sheet_name}: {img_error}")
                            continue

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from XLSX: {str(e)}',
                'images': []
            }


def extract_document_with_images(document_path: str, output_dir: str = None, markdown_dir: str = None, min_image_size: int = 50) -> Dict:
    """Extract complete document content with images inserted at their original positions"""
    try:
        extractor = ImageExtractor(document_path, output_dir, markdown_dir, min_image_size)
        result = extractor.extract_document_content_with_images()
        return result
    except Exception as e:
        return {
            'success': False,
            'error': f'Failed to extract document with images: {str(e)}',
            'markdown_content': '',
            'images': []
        }


def extract_images_from_document(document_path: str, output_dir: str = None, markdown_dir: str = None, markdown_mode: str = "simple", min_image_size: int = 50) -> Dict:
    """Main function to extract images from a document

    Args:
        document_path: Path to the document to extract images from
        output_dir: Directory to save extracted images
        markdown_dir: Directory where markdown files will be located (for relative path calculation)
        markdown_mode: How to format markdown references ("simple", "grouped", "inline")
    """
    try:
        extractor = ImageExtractor(document_path, output_dir, markdown_dir, min_image_size)

        result = extractor.extract_images()

        # Add different types of markdown references if images were extracted
        if result['success'] and result['images']:
            # Traditional markdown references (for backward compatibility)
            result['markdown_references'] = extractor.generate_markdown_references(result['images'], markdown_mode)

            # Additional formats for flexibility
            result['image_references_by_page'] = extractor.get_image_references_by_page(result['images'])
            result['simple_image_list'] = extractor.get_simple_image_list(result['images'])

            # Summary information
            result['images_by_page'] = {}
            for img in result['images']:
                page = img.get('page', img.get('slide', 'unknown'))
                if page not in result['images_by_page']:
                    result['images_by_page'][page] = 0
                result['images_by_page'][page] += 1
        else:
            result['markdown_references'] = ""
            result['image_references_by_page'] = {}
            result['simple_image_list'] = []
            result['images_by_page'] = {}

        return result

    except Exception as e:
        return {
            'success': False,
            'error': f'Failed to extract images: {str(e)}',
            'images': [],
            'markdown_references': "",
            'image_references_by_page': {},
            'simple_image_list': [],
            'images_by_page': {}
        }


def main():
    """Command line interface for image extraction"""
    if len(sys.argv) < 2:
        print("DocuGenius Image Extractor")
        print("Usage: python image_extractor.py <document_path> [output_dir] [markdown_dir] [mode] [min_image_size]")
        print("")
        print("Arguments:")
        print("  document_path  : Path to PDF, DOCX, or PPTX file")
        print("  output_dir     : Directory to save images (optional)")
        print("  markdown_dir   : Directory where markdown files will be located (optional)")
        print("  mode           : Extraction mode (optional)")
        print("                   - 'images_only': Extract images only (default)")
        print("                   - 'full_content': Extract full document with images in original positions")
        print("  min_image_size : Minimum image size in pixels (optional, default: 50)")
        print("")
        print("Output: JSON with extraction results")
        sys.exit(1)

    document_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    markdown_dir = sys.argv[3] if len(sys.argv) > 3 else None
    mode = sys.argv[4] if len(sys.argv) > 4 else "images_only"
    min_image_size = int(sys.argv[5]) if len(sys.argv) > 5 else 50

    if mode == "full_content":
        # Extract complete document content with images in original positions
        result = extract_document_with_images(document_path, output_dir, markdown_dir, min_image_size)
    else:
        # Extract images only (traditional mode)
        result = extract_images_from_document(document_path, output_dir, markdown_dir, "simple", min_image_size)

    # Output JSON result
    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == '__main__':
    main()
