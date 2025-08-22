# -*- coding: utf-8 -*-
import sys
import os
from pathlib import Path

# Ensure UTF-8 encoding on Windows
if sys.platform == 'win32':
    import codecs
    if hasattr(sys.stdout, 'reconfigure'):
        try:
            sys.stdout.reconfigure(encoding='utf-8')
            sys.stderr.reconfigure(encoding='utf-8')
        except:
            pass

def simple_convert(file_path):
    """Simple document converter with Windows compatibility"""
    try:
        # Handle Windows path issues
        file_path = os.path.normpath(file_path)
        if not os.path.exists(file_path):
            return f"Error: File not found: {file_path}"

        ext = Path(file_path).suffix.lower()
        name = Path(file_path).name

# Remove debug output to keep conversion clean
        
        if ext in ['.txt', '.md', '.markdown']:
            try:
                # Try UTF-8 first, then fallback to other encodings
                encodings = ['utf-8', 'utf-8-sig', 'gbk', 'cp1252', 'latin-1']
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            return f.read()
                    except UnicodeDecodeError:
                        continue
                return f"Error: Could not decode {name} with any supported encoding"
            except Exception as e:
                return f"Error reading {name}: {str(e)}"
                
        elif ext == '.docx':
            try:
                import docx
                doc = docx.Document(file_path)
                content = ""

                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        content += para.text.strip() + "\n\n"

                # Extract tables
                for table in doc.tables:
                    for i, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        if i == 0:  # Header row
                            content += "| " + " | ".join(row_data) + " |\n"
                            content += "| " + " | ".join(["---"] * len(row_data)) + " |\n"
                        else:
                            content += "| " + " | ".join(row_data) + " |\n"
                    content += "\n"

                return content.strip()
            except ImportError:
                return "Error: python-docx library not installed. Run: pip install python-docx"
            except Exception as e:
                return f"Error: {str(e)}"

        elif ext == '.xlsx':
            try:
                import openpyxl
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                content = ""

                for sheet_name in workbook.sheetnames:
                    if len(workbook.sheetnames) > 1:
                        content += f"## {sheet_name}\n\n"

                    worksheet = workbook[sheet_name]
                    rows = list(worksheet.iter_rows(values_only=True))

                    if rows:
                        # Filter out empty rows
                        non_empty_rows = []
                        for row in rows:
                            if any(cell is not None and str(cell).strip() for cell in row):
                                non_empty_rows.append(row)

                        if non_empty_rows:
                            # Find max columns
                            max_cols = max(len([cell for cell in row if cell is not None]) for row in non_empty_rows)

                            for i, row in enumerate(non_empty_rows):
                                row_data = []
                                for j in range(max_cols):
                                    if j < len(row) and row[j] is not None:
                                        row_data.append(str(row[j]).strip())
                                    else:
                                        row_data.append("")

                                if i == 0:  # Header row
                                    content += "| " + " | ".join(row_data) + " |\n"
                                    content += "| " + " | ".join(["---"] * len(row_data)) + " |\n"
                                else:
                                    content += "| " + " | ".join(row_data) + " |\n"
                            content += "\n"

                return content.strip()
            except ImportError:
                return "Error: openpyxl library not installed. Run: pip install openpyxl"
            except Exception as e:
                return f"Error: {str(e)}"

        elif ext == '.pptx':
            try:
                import pptx
                presentation = pptx.Presentation(file_path)
                content = ""

                for i, slide in enumerate(presentation.slides, 1):
                    if len(presentation.slides) > 1:
                        content += f"## Slide {i}\n\n"

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content += shape.text.strip() + "\n\n"

                    if i < len(presentation.slides):
                        content += "---\n\n"

                return content.strip()
            except ImportError:
                return "Error: python-pptx library not installed. Run: pip install python-pptx"
            except Exception as e:
                return f"Error: {str(e)}"

        elif ext == '.pdf':
            try:
                # Try pdfplumber first (better for Chinese text)
                try:
                    import pdfplumber
                    content = ""
                    with pdfplumber.open(file_path) as pdf:
                        for i, page in enumerate(pdf.pages):
                            if len(pdf.pages) > 1:
                                content += f"## Page {i+1}\n\n"

                            text = page.extract_text()
                            if text and text.strip():
                                # Clean up the text
                                lines = text.split('\n')
                                cleaned_lines = []
                                for line in lines:
                                    line = line.strip()
                                    if line:
                                        cleaned_lines.append(line)
                                content += '\n'.join(cleaned_lines) + "\n\n"

                    return content.strip()
                except ImportError:
                    raise ImportError("pdfplumber library not installed. Run: pip install pdfplumber")

            except ImportError:
                return "Error: pdfplumber library not installed. Run: pip install pdfplumber"
            except Exception as e:
                return f"Error: {str(e)}"
                
        else:
            return f"Unsupported file type: {ext}"
            
    except Exception as e:
        return f"Error: {str(e)}"

def convert_with_images(file_path, extract_images=True):
    """Convert document and optionally extract images"""
    try:
        # First, convert the document to markdown
        markdown_content = simple_convert(file_path)

        # If image extraction is enabled and we have a supported format
        if extract_images and Path(file_path).suffix.lower() in ['.pdf', '.docx', '.pptx', '.xlsx']:
            try:
                # Import and use the intelligent image extractor
                from image_extractor import extract_document_with_images

                # Extract complete document content with images in original positions
                extraction_result = extract_document_with_images(file_path)

                if extraction_result['success']:
                    if extraction_result.get('markdown_content'):
                        # Replace the markdown content with the intelligent version
                        # that has images in their original positions
                        markdown_content = extraction_result['markdown_content']

                        # Add metadata comment
                        if extraction_result.get('images_count', 0) > 0:
                            markdown_content += f"\n\n<!-- Images extracted: {extraction_result['images_count']} images saved to {extraction_result.get('output_dir', 'DocuGenius/images')} -->\n"
                    else:
                        # Fallback to traditional mode if no content was extracted
                        from image_extractor import extract_images_from_document
                        fallback_result = extract_images_from_document(file_path)
                        if fallback_result['success'] and fallback_result['images']:
                            markdown_content += fallback_result.get('markdown_references', '')
                            markdown_content += f"\n\n<!-- Images extracted: {fallback_result['images_count']} images saved to {fallback_result['output_dir']} -->\n"

            except ImportError:
                # Image extractor not available, continue without image extraction
                pass
            except Exception as e:
                # Image extraction failed, add note to markdown
                markdown_content += f"\n\n<!-- Note: Image extraction failed: {str(e)} -->\n"

        return markdown_content

    except Exception as e:
        return f"Error: {str(e)}"

if __name__ == '__main__':
    if len(sys.argv) > 1:
        try:
            # Default to extract images unless explicitly disabled
            extract_images = True
            if len(sys.argv) > 2 and sys.argv[2].lower() in ['false', 'no', '0']:
                extract_images = False

            result = convert_with_images(sys.argv[1], extract_images)
            if result:
                print(result)
            else:
                print("Error: Conversion failed")
                sys.exit(1)
        except Exception as e:
            print(f"Error: {str(e)}")
            sys.exit(1)
    else:
        print("Usage: converter.py file_path [extract_images=true]")
        print("  extract_images: true/false to enable/disable image extraction (default: true)")
        sys.exit(1)
