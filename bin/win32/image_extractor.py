# -*- coding: utf-8 -*-
"""
DocuGenius Image Extraction Module
Extracts images from PDF, DOCX, and PPTX files and organizes them in structured folders.
"""

import sys
import os
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Tuple, Optional

# Ensure UTF-8 encoding on Windows
if sys.platform == 'win32':
    import codecs
    if hasattr(sys.stdout, 'reconfigure'):
        try:
            sys.stdout.reconfigure(encoding='utf-8')
            sys.stderr.reconfigure(encoding='utf-8')
        except:
            pass

class ImageExtractor:
    """Main class for extracting images from documents"""

    def __init__(self, document_path: str, output_dir: str = None, markdown_dir: str = None, min_image_size: int = 50):
        self.document_path = Path(document_path)
        self.document_name = self.document_path.stem
        self.document_ext = self.document_path.suffix.lower()

        # Set output directory - use DocuGenius/images/{document_name}/ structure
        if output_dir:
            self.output_dir = Path(output_dir) / self.document_name
        else:
            # Default to DocuGenius/images folder in same directory as document
            self.output_dir = self.document_path.parent / "DocuGenius" / "images" / self.document_name

        self.extracted_images = []
        self.image_counter = 1
        self.min_image_size = min_image_size

        # Store the original document directory for relative path calculation
        self.original_dir = self.document_path.parent

        # Set markdown directory for relative path calculation
        if markdown_dir:
            self.markdown_dir = Path(markdown_dir)
        else:
            # Default to DocuGenius directory for markdown files
            self.markdown_dir = self.original_dir / "DocuGenius"
        
    def extract_images(self) -> Dict:
        """Extract images from document and return metadata"""
        try:
            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)
            
            if self.document_ext == '.pdf':
                return self._extract_from_pdf()
            elif self.document_ext == '.docx':
                return self._extract_from_docx()
            elif self.document_ext == '.pptx':
                return self._extract_from_pptx()
            elif self.document_ext == '.xlsx':
                return self._extract_from_xlsx()
            else:
                return {
                    'success': False,
                    'error': f'Unsupported document type: {self.document_ext}',
                    'images': []
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images: {str(e)}',
                'images': []
            }
    
    def _extract_from_pdf(self) -> Dict:
        """Extract images from PDF using PyMuPDF (fitz)"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            try:
                # Fallback to pdfplumber for basic image detection
                import pdfplumber
                return self._extract_from_pdf_pdfplumber()
            except ImportError:
                return {
                    'success': False,
                    'error': 'PDF image extraction requires PyMuPDF (pip install PyMuPDF) or pdfplumber',
                    'images': []
                }
        
        try:
            doc = fitz.open(str(self.document_path))
            images_extracted = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    # Skip if image is too small (likely decorative)
                    if pix.width < self.min_image_size or pix.height < self.min_image_size:
                        pix = None
                        continue
                    
                    # Determine image format and extension
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_ext = "png"
                        img_data = pix.tobytes("png")
                    else:  # CMYK: convert to RGB first
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_ext = "png"
                        img_data = pix1.tobytes("png")
                        pix1 = None
                    
                    # Generate unique filename
                    img_filename = self._generate_image_filename(
                        f"page_{page_num + 1}_img_{img_index + 1}", 
                        img_ext
                    )
                    img_path = self.output_dir / img_filename
                    
                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(img_data)
                    
                    # Add to extracted images list
                    image_info = {
                        'filename': img_filename,
                        'path': str(img_path),
                        'relative_path': self._calculate_relative_path(img_path),
                        'page': page_num + 1,
                        'width': pix.width,
                        'height': pix.height,
                        'format': img_ext.upper(),
                        'size_bytes': len(img_data)
                    }
                    images_extracted.append(image_info)
                    self.extracted_images.append(image_info)
                    
                    pix = None
            
            doc.close()
            
            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from PDF: {str(e)}',
                'images': []
            }
    
    def _extract_from_pdf_pdfplumber(self) -> Dict:
        """Fallback PDF image extraction using pdfplumber (limited functionality)"""
        try:
            import pdfplumber
            
            # pdfplumber doesn't directly extract images, but we can detect them
            # This is a basic implementation that creates placeholder entries
            with pdfplumber.open(str(self.document_path)) as pdf:
                images_detected = []
                
                for page_num, page in enumerate(pdf.pages):
                    # Check if page has images (basic detection)
                    if hasattr(page, 'images') and page.images:
                        for img_index, img_obj in enumerate(page.images):
                            # Create placeholder entry since we can't extract actual image data
                            image_info = {
                                'filename': f"page_{page_num + 1}_img_{img_index + 1}_placeholder.txt",
                                'path': str(self.output_dir / f"page_{page_num + 1}_img_{img_index + 1}_placeholder.txt"),
                                'relative_path': f"images/{self.document_name}/page_{page_num + 1}_img_{img_index + 1}_placeholder.txt",
                                'page': page_num + 1,
                                'width': img_obj.get('width', 0),
                                'height': img_obj.get('height', 0),
                                'format': 'PLACEHOLDER',
                                'size_bytes': 0,
                                'note': 'Image detected but not extracted (requires PyMuPDF for full extraction)'
                            }
                            images_detected.append(image_info)
                            
                            # Create placeholder file
                            placeholder_path = self.output_dir / image_info['filename']
                            with open(placeholder_path, 'w') as f:
                                f.write(f"Image placeholder for page {page_num + 1}, image {img_index + 1}\n")
                                f.write("Install PyMuPDF (pip install PyMuPDF) for actual image extraction.\n")
                
                return {
                    'success': True,
                    'document': str(self.document_path),
                    'output_dir': str(self.output_dir),
                    'images_count': len(images_detected),
                    'images': images_detected,
                    'note': 'Limited extraction - install PyMuPDF for full image extraction'
                }
                
        except Exception as e:
            return {
                'success': False,
                'error': f'Error with pdfplumber fallback: {str(e)}',
                'images': []
            }

    def _extract_from_docx(self) -> Dict:
        """Extract images from DOCX files"""
        try:
            import docx
            from docx.document import Document
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            return {
                'success': False,
                'error': 'DOCX image extraction requires python-docx (pip install python-docx)',
                'images': []
            }

        try:
            doc = docx.Document(str(self.document_path))
            images_extracted = []

            # Extract images from document relationships
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    # Get image data
                    image_part = rel.target_part
                    image_data = image_part.blob

                    # Determine file extension from content type
                    content_type = image_part.content_type
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        img_ext = 'jpg'
                    elif 'png' in content_type:
                        img_ext = 'png'
                    elif 'gif' in content_type:
                        img_ext = 'gif'
                    elif 'bmp' in content_type:
                        img_ext = 'bmp'
                    else:
                        img_ext = 'png'  # Default fallback

                    # Generate unique filename
                    img_filename = self._generate_image_filename(
                        f"docx_img_{len(images_extracted) + 1}",
                        img_ext
                    )
                    img_path = self.output_dir / img_filename

                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_data)

                    # Add to extracted images list
                    image_info = {
                        'filename': img_filename,
                        'path': str(img_path),
                        'relative_path': self._calculate_relative_path(img_path),
                        'format': img_ext.upper(),
                        'size_bytes': len(image_data),
                        'source': 'docx_relationship'
                    }
                    images_extracted.append(image_info)
                    self.extracted_images.append(image_info)

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from DOCX: {str(e)}',
                'images': []
            }

    def _extract_from_pptx(self) -> Dict:
        """Extract images from PPTX files"""
        try:
            import pptx
            from pptx import Presentation
        except ImportError:
            return {
                'success': False,
                'error': 'PPTX image extraction requires python-pptx (pip install python-pptx)',
                'images': []
            }

        try:
            prs = Presentation(str(self.document_path))
            images_extracted = []

            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        # Get image data
                        image = shape.image
                        image_data = image.blob

                        # Determine file extension from content type
                        content_type = image.content_type
                        if 'jpeg' in content_type or 'jpg' in content_type:
                            img_ext = 'jpg'
                        elif 'png' in content_type:
                            img_ext = 'png'
                        elif 'gif' in content_type:
                            img_ext = 'gif'
                        elif 'bmp' in content_type:
                            img_ext = 'bmp'
                        else:
                            img_ext = 'png'  # Default fallback

                        # Generate unique filename
                        img_filename = self._generate_image_filename(
                            f"slide_{slide_num + 1}_img_{len([s for s in slide.shapes if hasattr(s, 'image')]) + 1}",
                            img_ext
                        )
                        img_path = self.output_dir / img_filename

                        # Save image
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_data)

                        # Add to extracted images list
                        image_info = {
                            'filename': img_filename,
                            'path': str(img_path),
                            'relative_path': self._calculate_relative_path(img_path),
                            'slide': slide_num + 1,
                            'format': img_ext.upper(),
                            'size_bytes': len(image_data),
                            'source': 'pptx_shape'
                        }
                        images_extracted.append(image_info)
                        self.extracted_images.append(image_info)

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from PPTX: {str(e)}',
                'images': []
            }

    def _extract_from_xlsx(self) -> Dict:
        """Extract images from XLSX files"""
        try:
            import openpyxl
            from openpyxl.drawing.image import Image as OpenpyxlImage
        except ImportError:
            return {
                'success': False,
                'error': 'XLSX image extraction requires openpyxl (pip install openpyxl)',
                'images': []
            }

        try:
            workbook = openpyxl.load_workbook(str(self.document_path))
            images_extracted = []

            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]

                # Check if worksheet has images
                if hasattr(worksheet, '_images') and worksheet._images:
                    for img_index, image in enumerate(worksheet._images):
                        try:
                            # Get image data
                            image_data = image.ref.getvalue()

                            # Determine file extension from image format
                            img_ext = 'png'  # Default
                            if hasattr(image, 'format') and image.format:
                                img_ext = image.format.lower()
                            elif hasattr(image.ref, 'format') and image.ref.format:
                                img_ext = image.ref.format.lower()

                            # Generate unique filename
                            img_filename = self._generate_image_filename(
                                f"sheet_{sheet_name}_img_{img_index + 1}",
                                img_ext
                            )
                            img_path = self.output_dir / img_filename

                            # Save image
                            with open(img_path, "wb") as img_file:
                                img_file.write(image_data)

                            # Add to extracted images list
                            image_info = {
                                'filename': img_filename,
                                'path': str(img_path),
                                'relative_path': self._calculate_relative_path(img_path),
                                'sheet': sheet_name,
                                'format': img_ext.upper(),
                                'size_bytes': len(image_data),
                                'source': 'xlsx_image'
                            }
                            images_extracted.append(image_info)
                            self.extracted_images.append(image_info)

                        except Exception as img_error:
                            # Skip this image if there's an error
                            print(f"Warning: Could not extract image {img_index + 1} from sheet {sheet_name}: {img_error}")
                            continue

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(images_extracted),
                'images': images_extracted
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting images from XLSX: {str(e)}',
                'images': []
            }

    def _generate_image_filename(self, base_name: str, extension: str) -> str:
        """Generate a unique filename for an image, avoiding collisions"""
        # Clean the base name
        base_name = "".join(c for c in base_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
        base_name = base_name.replace(' ', '_')

        filename = f"{base_name}.{extension}"
        counter = 1

        # Check for collisions and add counter if needed
        while (self.output_dir / filename).exists():
            filename = f"{base_name}_{counter}.{extension}"
            counter += 1

        return filename

    def _calculate_relative_path(self, image_path: Path) -> str:
        """Calculate the correct relative path from markdown file to image"""
        try:
            # Use the instance markdown_dir for relative path calculation
            markdown_dir = self.markdown_dir

            # Calculate relative path from markdown file to image
            relative_path = os.path.relpath(str(image_path), str(markdown_dir))

            # Ensure forward slashes for markdown compatibility
            relative_path = relative_path.replace('\\', '/')

            return relative_path
        except Exception as e:
            # Fallback to simple relative path
            return f"images/{self.document_name}/{image_path.name}"

    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate MD5 hash of a file for duplicate detection"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def generate_markdown_references(self, images: List[Dict], mode: str = "simple") -> str:
        """Generate markdown image references for extracted images

        Args:
            images: List of image dictionaries with metadata
            mode: "simple" - just image paths, "grouped" - grouped by page/slide, "inline" - individual references
        """
        if not images:
            return ""

        if mode == "simple":
            # Just return the image paths as simple markdown references
            markdown_lines = []
            for img in images:
                alt_text = f"Image"
                if 'page' in img:
                    alt_text += f" from page {img['page']}"
                elif 'slide' in img:
                    alt_text += f" from slide {img['slide']}"

                markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
            return "\n\n".join(markdown_lines)

        elif mode == "grouped":
            # Group images by page/slide
            grouped = {}
            for img in images:
                key = img.get('page', img.get('slide', 'unknown'))
                if key not in grouped:
                    grouped[key] = []
                grouped[key].append(img)

            markdown_lines = []
            for key in sorted(grouped.keys()):
                if isinstance(key, int):
                    if 'page' in grouped[key][0]:
                        markdown_lines.append(f"\n### Page {key} Images\n")
                    else:
                        markdown_lines.append(f"\n### Slide {key} Images\n")
                else:
                    markdown_lines.append(f"\n### Images\n")

                for img in grouped[key]:
                    alt_text = f"Image from {img.get('source', 'document')}"
                    markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
                    markdown_lines.append("")

            return "\n".join(markdown_lines)

        else:  # "inline" or default
            # Traditional mode - all images at the end
            markdown_lines = []
            markdown_lines.append("\n## Extracted Images\n")

            for img in images:
                alt_text = f"Image from {img.get('source', 'document')}"
                if 'page' in img:
                    alt_text += f" (Page {img['page']})"
                elif 'slide' in img:
                    alt_text += f" (Slide {img['slide']})"

                markdown_lines.append(f"![{alt_text}]({img['relative_path']})")
                markdown_lines.append("")  # Empty line for spacing

            return "\n".join(markdown_lines)

    def get_image_references_by_page(self, images: List[Dict]) -> Dict[int, List[str]]:
        """Get image references organized by page number for easy insertion"""
        page_images = {}

        for img in images:
            page = img.get('page', img.get('slide', 0))
            if page not in page_images:
                page_images[page] = []

            alt_text = f"Image"
            if 'page' in img:
                alt_text += f" from page {img['page']}"
            elif 'slide' in img:
                alt_text += f" from slide {img['slide']}"

            page_images[page].append(f"![{alt_text}]({img['relative_path']})")

        return page_images

    def get_simple_image_list(self, images: List[Dict]) -> List[str]:
        """Get a simple list of image markdown references"""
        image_refs = []

        for img in images:
            alt_text = f"Image"
            if 'page' in img:
                alt_text += f" from page {img['page']}"
            elif 'slide' in img:
                alt_text += f" from slide {img['slide']}"

            image_refs.append(f"![{alt_text}]({img['relative_path']})")

        return image_refs

    def extract_document_content_with_images(self) -> Dict:
        """Extract document content and intelligently insert images at their original positions"""
        try:
            if self.document_ext == '.pdf':
                return self._extract_pdf_content_with_images()
            elif self.document_ext == '.docx':
                return self._extract_docx_content_with_images()
            elif self.document_ext == '.pptx':
                return self._extract_pptx_content_with_images()
            elif self.document_ext == '.xlsx':
                return self._extract_xlsx_content_with_images()
            else:
                return {
                    'success': False,
                    'error': f'Unsupported document type for content extraction: {self.document_ext}',
                    'markdown_content': '',
                    'images': []
                }
        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_pdf_content_with_images(self) -> Dict:
        """Extract PDF content and insert images at their original positions"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {
                'success': False,
                'error': 'PDF content extraction requires PyMuPDF (pip install PyMuPDF)',
                'markdown_content': '',
                'images': []
            }

        try:
            doc = fitz.open(str(self.document_path))
            markdown_lines = []
            all_images = []

            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)

                # Add page header
                markdown_lines.append(f"\n## Page {page_num + 1}\n")

                # Extract text content
                text_content = page.get_text()
                if text_content.strip():
                    # Split into paragraphs and clean up
                    paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
                    for paragraph in paragraphs:
                        if paragraph:
                            markdown_lines.append(paragraph)
                            markdown_lines.append("")  # Empty line after paragraph

                # Extract and insert images for this page
                image_list = page.get_images()
                page_images = []

                for img_index, img in enumerate(image_list):
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)

                    # Skip if image is too small (likely decorative)
                    if pix.width < self.min_image_size or pix.height < self.min_image_size:
                        pix = None
                        continue

                    # Determine image format and extension
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_ext = "png"
                        img_data = pix.tobytes("png")
                    else:  # CMYK: convert to RGB first
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_ext = "png"
                        img_data = pix1.tobytes("png")
                        pix1 = None

                    # Generate unique filename
                    img_filename = self._generate_image_filename(
                        f"page_{page_num + 1}_img_{img_index + 1}",
                        img_ext
                    )
                    img_path = self.output_dir / img_filename

                    # Save image
                    with open(img_path, "wb") as img_file:
                        img_file.write(img_data)

                    # Create image info
                    image_info = {
                        'filename': img_filename,
                        'path': str(img_path),
                        'relative_path': self._calculate_relative_path(img_path),
                        'page': page_num + 1,
                        'width': pix.width,
                        'height': pix.height,
                        'format': img_ext.upper(),
                        'size_bytes': len(img_data)
                    }
                    page_images.append(image_info)
                    all_images.append(image_info)

                    pix = None

                # Insert images after the text content for this page
                if page_images:
                    for img_info in page_images:
                        alt_text = f"Image from page {img_info['page']}"
                        markdown_lines.append(f"![{alt_text}]({img_info['relative_path']})")
                        markdown_lines.append("")  # Empty line after image

                markdown_lines.append("---\n")  # Page separator

            doc.close()

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(all_images),
                'images': all_images,
                'markdown_content': '\n'.join(markdown_lines)
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting PDF content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_docx_content_with_images(self) -> Dict:
        """Extract DOCX content and insert images at their original positions"""
        try:
            import docx
            from docx.document import Document
        except ImportError:
            return {
                'success': False,
                'error': 'DOCX content extraction requires python-docx (pip install python-docx)',
                'markdown_content': '',
                'images': []
            }

        try:
            doc = docx.Document(str(self.document_path))
            markdown_lines = []
            all_images = []

            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            # Extract images from document relationships first
            image_parts = {}
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_parts[rel.rId] = rel.target_part

            img_counter = 1

            # Process document elements in order
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    paragraph = docx.text.paragraph.Paragraph(element, doc)
                    text = paragraph.text.strip()

                    if text:
                        # Check if paragraph has special formatting
                        if paragraph.style.name.startswith('Heading'):
                            level = '##' if 'Heading 1' in paragraph.style.name else '###'
                            markdown_lines.append(f"{level} {text}\n")
                        else:
                            markdown_lines.append(text)
                            markdown_lines.append("")  # Empty line after paragraph

                    # Check for images in this paragraph
                    for run in paragraph.runs:
                        for drawing in run.element.xpath('.//a:blip'):
                            rId = drawing.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if rId in image_parts:
                                image_part = image_parts[rId]
                                image_data = image_part.blob

                                # Determine file extension from content type
                                content_type = image_part.content_type
                                if 'jpeg' in content_type or 'jpg' in content_type:
                                    img_ext = 'jpg'
                                elif 'png' in content_type:
                                    img_ext = 'png'
                                elif 'gif' in content_type:
                                    img_ext = 'gif'
                                else:
                                    img_ext = 'png'  # Default fallback

                                # Generate unique filename
                                img_filename = self._generate_image_filename(
                                    f"docx_img_{img_counter}",
                                    img_ext
                                )
                                img_path = self.output_dir / img_filename

                                # Save image
                                with open(img_path, "wb") as img_file:
                                    img_file.write(image_data)

                                # Create image info
                                image_info = {
                                    'filename': img_filename,
                                    'path': str(img_path),
                                    'relative_path': self._calculate_relative_path(img_path),
                                    'format': img_ext.upper(),
                                    'size_bytes': len(image_data),
                                    'source': 'docx_inline'
                                }
                                all_images.append(image_info)

                                # Insert image reference right here in the content
                                alt_text = f"Image {img_counter}"
                                markdown_lines.append(f"![{alt_text}]({image_info['relative_path']})")
                                markdown_lines.append("")  # Empty line after image

                                img_counter += 1

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(all_images),
                'images': all_images,
                'markdown_content': '\n'.join(markdown_lines)
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting DOCX content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_pptx_content_with_images(self) -> Dict:
        """Extract PPTX content and insert images at their original positions"""
        try:
            import pptx
            from pptx import Presentation
        except ImportError:
            return {
                'success': False,
                'error': 'PPTX content extraction requires python-pptx (pip install python-pptx)',
                'markdown_content': '',
                'images': []
            }

        try:
            prs = Presentation(str(self.document_path))
            markdown_lines = []
            all_images = []

            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            for slide_num, slide in enumerate(prs.slides):
                # Add slide header
                markdown_lines.append(f"\n## Slide {slide_num + 1}\n")

                # Extract text and images from shapes in order
                for shape in slide.shapes:
                    # Handle text shapes
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        if text:
                            markdown_lines.append(text)
                            markdown_lines.append("")  # Empty line after text

                    # Handle image shapes
                    if hasattr(shape, 'image'):
                        image = shape.image
                        image_data = image.blob

                        # Determine file extension from content type
                        content_type = image.content_type
                        if 'jpeg' in content_type or 'jpg' in content_type:
                            img_ext = 'jpg'
                        elif 'png' in content_type:
                            img_ext = 'png'
                        elif 'gif' in content_type:
                            img_ext = 'gif'
                        else:
                            img_ext = 'png'  # Default fallback

                        # Generate unique filename
                        img_filename = self._generate_image_filename(
                            f"slide_{slide_num + 1}_img_{len([s for s in slide.shapes if hasattr(s, 'image')]) + 1}",
                            img_ext
                        )
                        img_path = self.output_dir / img_filename

                        # Save image
                        with open(img_path, "wb") as img_file:
                            img_file.write(image_data)

                        # Create image info
                        image_info = {
                            'filename': img_filename,
                            'path': str(img_path),
                            'relative_path': self._calculate_relative_path(img_path),
                            'slide': slide_num + 1,
                            'format': img_ext.upper(),
                            'size_bytes': len(image_data),
                            'source': 'pptx_shape'
                        }
                        all_images.append(image_info)

                        # Insert image reference right here in the content
                        alt_text = f"Image from slide {slide_num + 1}"
                        markdown_lines.append(f"![{alt_text}]({image_info['relative_path']})")
                        markdown_lines.append("")  # Empty line after image

                markdown_lines.append("---\n")  # Slide separator

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(all_images),
                'images': all_images,
                'markdown_content': '\n'.join(markdown_lines)
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting PPTX content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }

    def _extract_xlsx_content_with_images(self) -> Dict:
        """Extract XLSX content and insert images at their original positions"""
        try:
            import openpyxl
        except ImportError:
            return {
                'success': False,
                'error': 'XLSX content extraction requires openpyxl (pip install openpyxl)',
                'markdown_content': '',
                'images': []
            }

        try:
            workbook = openpyxl.load_workbook(str(self.document_path))
            markdown_lines = []
            all_images = []

            # Create output directory
            self.output_dir.mkdir(parents=True, exist_ok=True)

            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]

                # Add sheet header
                markdown_lines.append(f"\n## {sheet_name}\n")

                # Extract table data
                rows = list(worksheet.iter_rows(values_only=True))
                if rows:
                    # Filter out completely empty rows
                    non_empty_rows = []
                    for row in rows:
                        if any(cell is not None and str(cell).strip() for cell in row):
                            non_empty_rows.append(row)

                    if non_empty_rows:
                        # Create markdown table
                        markdown_lines.append("| " + " | ".join(str(cell) if cell is not None else "" for cell in non_empty_rows[0]) + " |")
                        markdown_lines.append("| " + " | ".join("---" for _ in non_empty_rows[0]) + " |")

                        for row in non_empty_rows[1:]:
                            markdown_lines.append("| " + " | ".join(str(cell) if cell is not None else "" for cell in row) + " |")

                        markdown_lines.append("")  # Empty line after table

                # Extract and insert images for this sheet
                if hasattr(worksheet, '_images') and worksheet._images:
                    sheet_images = []

                    for img_index, image in enumerate(worksheet._images):
                        try:
                            # Get image data
                            image_data = image.ref.getvalue()

                            # Determine file extension
                            img_ext = 'png'  # Default
                            if hasattr(image, 'format') and image.format:
                                img_ext = image.format.lower()
                            elif hasattr(image.ref, 'format') and image.ref.format:
                                img_ext = image.ref.format.lower()

                            # Generate unique filename
                            img_filename = self._generate_image_filename(
                                f"sheet_{sheet_name}_img_{img_index + 1}",
                                img_ext
                            )
                            img_path = self.output_dir / img_filename

                            # Save image
                            with open(img_path, "wb") as img_file:
                                img_file.write(image_data)

                            # Create image info
                            image_info = {
                                'filename': img_filename,
                                'path': str(img_path),
                                'relative_path': self._calculate_relative_path(img_path),
                                'sheet': sheet_name,
                                'format': img_ext.upper(),
                                'size_bytes': len(image_data),
                                'source': 'xlsx_image'
                            }
                            sheet_images.append(image_info)
                            all_images.append(image_info)

                        except Exception as img_error:
                            print(f"Warning: Could not extract image {img_index + 1} from sheet {sheet_name}: {img_error}")
                            continue

                    # Insert images after the sheet content
                    if sheet_images:
                        for img_info in sheet_images:
                            alt_text = f"Image from sheet {img_info['sheet']}"
                            markdown_lines.append(f"![{alt_text}]({img_info['relative_path']})")
                            markdown_lines.append("")  # Empty line after image

                markdown_lines.append("---\n")  # Sheet separator

            return {
                'success': True,
                'document': str(self.document_path),
                'output_dir': str(self.output_dir),
                'images_count': len(all_images),
                'images': all_images,
                'markdown_content': '\n'.join(markdown_lines)
            }

        except Exception as e:
            return {
                'success': False,
                'error': f'Error extracting XLSX content with images: {str(e)}',
                'markdown_content': '',
                'images': []
            }


def extract_images_from_document(document_path: str, output_dir: str = None, markdown_dir: str = None, markdown_mode: str = "simple", min_image_size: int = 50) -> Dict:
    """Main function to extract images from a document

    Args:
        document_path: Path to the document to extract images from
        output_dir: Directory to save extracted images
        markdown_dir: Directory where markdown files will be located (for relative path calculation)
        markdown_mode: How to format markdown references ("simple", "grouped", "inline")
    """
    try:
        extractor = ImageExtractor(document_path, output_dir, markdown_dir, min_image_size)

        result = extractor.extract_images()

        # Add different types of markdown references if images were extracted
        if result['success'] and result['images']:
            # Traditional markdown references (for backward compatibility)
            result['markdown_references'] = extractor.generate_markdown_references(result['images'], markdown_mode)

            # Additional formats for flexibility
            result['image_references_by_page'] = extractor.get_image_references_by_page(result['images'])
            result['simple_image_list'] = extractor.get_simple_image_list(result['images'])

            # Summary information
            result['images_by_page'] = {}
            for img in result['images']:
                page = img.get('page', img.get('slide', 'unknown'))
                if page not in result['images_by_page']:
                    result['images_by_page'][page] = 0
                result['images_by_page'][page] += 1
        else:
            result['markdown_references'] = ""
            result['image_references_by_page'] = {}
            result['simple_image_list'] = []
            result['images_by_page'] = {}

        return result

    except Exception as e:
        return {
            'success': False,
            'error': f'Failed to extract images: {str(e)}',
            'images': [],
            'markdown_references': "",
            'image_references_by_page': {},
            'simple_image_list': [],
            'images_by_page': {}
        }


def extract_document_with_images(document_path: str, output_dir: str = None, markdown_dir: str = None, min_image_size: int = 50) -> Dict:
    """Extract complete document content with images inserted at their original positions"""
    try:
        extractor = ImageExtractor(document_path, output_dir, markdown_dir, min_image_size)
        result = extractor.extract_document_content_with_images()
        return result
    except Exception as e:
        return {
            'success': False,
            'error': f'Failed to extract document with images: {str(e)}',
            'markdown_content': '',
            'images': []
        }


def main():
    """Command line interface for image extraction"""
    if len(sys.argv) < 2:
        print("DocuGenius Image Extractor")
        print("Usage: python image_extractor.py <document_path> [output_dir] [markdown_dir] [mode] [min_image_size]")
        print("")
        print("Arguments:")
        print("  document_path  : Path to PDF, DOCX, or PPTX file")
        print("  output_dir     : Directory to save images (optional)")
        print("  markdown_dir   : Directory where markdown files will be located (optional)")
        print("  mode           : Extraction mode (optional)")
        print("                   - 'images_only': Extract images only (default)")
        print("                   - 'full_content': Extract full document with images in original positions")
        print("  min_image_size : Minimum image size in pixels (optional, default: 50)")
        print("")
        print("Output: JSON with extraction results")
        sys.exit(1)

    document_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    markdown_dir = sys.argv[3] if len(sys.argv) > 3 else None
    mode = sys.argv[4] if len(sys.argv) > 4 else "images_only"
    min_image_size = int(sys.argv[5]) if len(sys.argv) > 5 else 50

    if mode == "full_content":
        # Extract complete document content with images in original positions
        result = extract_document_with_images(document_path, output_dir, markdown_dir, min_image_size)
    else:
        # Extract images only (traditional mode)
        result = extract_images_from_document(document_path, output_dir, markdown_dir, "simple", min_image_size)

    # Output JSON result
    print(json.dumps(result, indent=2, ensure_ascii=False))


if __name__ == '__main__':
    main()
